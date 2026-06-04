#!/usr/bin/env python3
"""Cross-platform PPTX structure validator for NEPU PPT output.

Performs structural checks on a PPTX file using python-pptx only.
No Office/WPS/LibreOffice required. Complements the visual QA process
in references/visual-qa.md.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Any


try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
except ImportError:
    print("python-pptx is required: pip install python-pptx", file=sys.stderr)
    raise SystemExit(2)


# Approximate character widths (pt) for common scripts
# CJK characters ≈ 1.0 em; Latin ≈ 0.55 em
_APPROX_CHAR_WIDTH = {
    "cjk": 1.0,   # multiplier of font size
    "latin": 0.55,
}

PLACEHOLDER_PATTERNS = [
    re.compile(r"\blorem\b", re.I),
    re.compile(r"\bipsum\b", re.I),
    re.compile(r"x{3,}"),
    re.compile(r"\btodo\b", re.I),
    re.compile(r"\bplaceholder\b", re.I),
    re.compile(r"\[TBD\]", re.I),
    re.compile(r"此处", re.I),
]


def slide_dims(prs: Any) -> tuple[float, float]:
    """Return slide width and height in inches."""
    return (prs.slide_width / 914400, prs.slide_height / 914400)


def _has_cjk(text: str) -> bool:
    return bool(re.search(r"[\u4e00-\u9fff\u3400-\u4dbf]", text))


def estimate_text_width_pt(text: str, font_size_pt: float) -> float:
    """Estimate rendered text width in points."""
    cjk_count = sum(1 for ch in text if _has_cjk(ch))
    latin_count = len(text) - cjk_count
    return (cjk_count * _APPROX_CHAR_WIDTH["cjk"] + latin_count * _APPROX_CHAR_WIDTH["latin"]) * font_size_pt


def check_shape_bounds(left: int, top: int, width: int, height: int,
                       slide_w: float, slide_h: float) -> list[str]:
    """Check a shape fits within slide canvas. Returns issue strings."""
    issues: list[str] = []
    right_emu = left + width
    bottom_emu = top + height
    slide_w_emu = int(slide_w * 914400)
    slide_h_emu = int(slide_h * 914400)

    if left < 0:
        issues.append(f"shape extends left of slide (left={left})")
    if top < 0:
        issues.append(f"shape extends above slide (top={top})")
    if right_emu > slide_w_emu:
        issues.append(f"shape extends right of slide (right={right_emu}, max={slide_w_emu})")
    if bottom_emu > slide_h_emu:
        issues.append(f"shape extends below slide (bottom={bottom_emu}, max={slide_h_emu})")
    return issues


def validate_pptx(path: Path) -> dict[str, Any]:
    """Run all structural checks and return a report dict."""

    report: dict[str, Any] = {
        "file": str(path),
        "valid": True,
        "errors": [],
        "warnings": [],
        "stats": {},
    }

    if not path.exists():
        report["valid"] = False
        report["errors"].append("File does not exist")
        return report

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        report["valid"] = False
        report["errors"].append(f"Cannot open as PPTX: {exc}")
        return report

    sw, sh = slide_dims(prs)
    slide_count = len(prs.slides)
    stats: dict[str, Any] = {
        "slide_count": slide_count,
        "slide_width_inches": round(sw, 2),
        "slide_height_inches": round(sh, 2),
        "aspect_ratio": round(sw / sh, 3) if sh else 0,
        "slides_with_notes": 0,
        "slides_with_images": 0,
        "text_overflow_slides": set(),
        "small_images": 0,
        "placeholder_matches": [],
        "total_shapes": 0,
        "total_text_boxes": 0,
    }

    # Per-slide checks
    for slide_idx, slide in enumerate(prs.slides, start=1):
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            stats["slides_with_notes"] += 1

        has_image_this_slide = False
        shape_types: list[str] = []

        for shape in slide.shapes:
            stats["total_shapes"] += 1
            shape_types.append(shape.shape_type)

            # Bounds check
            left = shape.left if shape.left is not None else 0
            top = shape.top if shape.top is not None else 0
            width = shape.width if shape.width is not None else 0
            height = shape.height if shape.height is not None else 0
            bound_issues = check_shape_bounds(left, top, width, height, sw, sh)
            for issue in bound_issues:
                report["warnings"].append(f"Slide {slide_idx}: {issue}")

            # Image checks
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                has_image_this_slide = True
                img_w_inches = shape.width / 914400 if shape.width else 0
                img_h_inches = shape.height / 914400 if shape.height else 0
                if img_w_inches < 1.0 and img_h_inches < 1.0:
                    stats["small_images"] += 1
                    report["warnings"].append(
                        f"Slide {slide_idx}: image '{shape.name}' is small "
                        f"({img_w_inches:.1f}\"×{img_h_inches:.1f}\"), may be unreadable"
                    )

            # Text overflow check
            if shape.has_text_frame:
                stats["total_text_boxes"] += 1
                for paragraph in shape.text_frame.paragraphs:
                    full_text = paragraph.text
                    if not full_text:
                        continue
                    # Placeholder scan
                    for pattern in PLACEHOLDER_PATTERNS:
                        if pattern.search(full_text):
                            report["warnings"].append(
                                f"Slide {slide_idx}: placeholder text found: '{full_text[:60]}'"
                            )
                            stats["placeholder_matches"].append(full_text[:60])

                    # Text overflow estimation
                    font_size = 12  # default
                    for run in paragraph.runs:
                        if run.font.size:
                            font_size = run.font.size / 12700  # EMU to pt
                            break
                    estimated_width = estimate_text_width_pt(full_text, font_size)
                    if shape.width:
                        shape_width_pt = shape.width / 12700
                        if estimated_width > shape_width_pt * 1.15:
                            stats["text_overflow_slides"].add(slide_idx)
                            break  # one warning per shape

        if has_image_this_slide:
            stats["slides_with_images"] += 1

        # Layout repeat detection: collect shape type signatures
        stats.setdefault("layout_signatures", []).append(tuple(sorted(shape_types)))

    # Post-process stats
    stats["text_overflow_slides"] = sorted(stats["text_overflow_slides"])
    stats["unique_layout_count"] = len(set(stats.get("layout_signatures", [])))

    # Layout repeat warning
    sigs = stats.get("layout_signatures", [])
    if sigs:
        from collections import Counter
        sig_counts = Counter(sigs)
        repeats = {k: v for k, v in sig_counts.items() if v >= 5}
        if repeats:
            report["warnings"].append(
                f"Layout pattern repeated ≥5 times in {dict(repeats)}. "
                f"Unique layouts: {stats['unique_layout_count']}/{stats['slide_count']}"
            )

    # Slide count
    if slide_count == 0:
        report["valid"] = False
        report["errors"].append("PPTX has zero slides")
    elif slide_count < 3:
        report["warnings"].append(f"Only {slide_count} slides — confirm this is intended")

    # Notes coverage
    if stats["slides_with_notes"] > 0 and stats["slides_with_notes"] < slide_count:
        report["warnings"].append(
            f"Only {stats['slides_with_notes']}/{slide_count} slides have speaker notes"
        )

    report["stats"] = {k: (sorted(v) if isinstance(v, set) else v)
                       for k, v in stats.items()}
    return report


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx", type=Path, help="Path to PPTX file to validate")
    parser.add_argument("--json", action="store_true", help="Output as JSON")
    args = parser.parse_args()

    report = validate_pptx(args.pptx)

    if args.json:
        print(json.dumps(report, ensure_ascii=False, indent=2, default=str))
    else:
        print(f"File: {report['file']}")
        print(f"Valid: {report['valid']}")
        stats = report.get("stats", {})
        print(f"Slides: {stats.get('slide_count', '?')} | "
              f"Size: {stats.get('slide_width_inches', '?')}\"×{stats.get('slide_height_inches', '?')}\" | "
              f"Aspect: {stats.get('aspect_ratio', '?')}")
        print(f"Shapes: {stats.get('total_shapes', 0)} | "
              f"Text boxes: {stats.get('total_text_boxes', 0)}")
        print(f"Slides with notes: {stats.get('slides_with_notes', 0)} | "
              f"Slides with images: {stats.get('slides_with_images', 0)}")
        print(f"Unique layouts: {stats.get('unique_layout_count', 0)}")
        print(f"Text overflow slides: {stats.get('text_overflow_slides', [])}")
        print(f"Placeholder matches: {len(stats.get('placeholder_matches', []))}")
        print(f"Small images: {stats.get('small_images', 0)}")

        errors = report.get("errors", [])
        warnings = report.get("warnings", [])
        if errors:
            print(f"\n❌ ERRORS ({len(errors)}):")
            for e in errors:
                print(f"  • {e}")
        if warnings:
            print(f"\n⚠ WARNINGS ({len(warnings)}):")
            for w in warnings[:20]:
                print(f"  • {w}")
            if len(warnings) > 20:
                print(f"  … and {len(warnings) - 20} more")

        if not errors and not warnings:
            print("\n✓ No issues found.")

    return 0 if report["valid"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
