"""High-level slide builder for NEPU-branded PowerPoint decks.

Provides reusable patterns (cover, section, figure, matrix, conclusion)
so AI agents produce consistent, editable PPTX without rewriting boilerplate
shape coordinates for every slide.

Usage:
    from scripts.slide_builder import NepuSlideBuilder

    builder = NepuSlideBuilder(theme="petro-blue")
    builder.add_cover("标题", subtitle="副标题")
    builder.add_figure_slide("figures/result.png", claim="核心发现")
    builder.add_section("第三章：实验结果")
    builder.add_matrix_slide(headers=["A","B"], rows=[["1","2"],["3","4"]])
    builder.add_conclusion(["结论一", "结论二"])
    builder.save("output.pptx")
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ── NEPU color palette ────────────────────────────────────────────
PALETTE: dict[str, str] = {
    "blue": "#00508F",
    "dark_blue": "#00335C",
    "red": "#C41230",
    "gold": "#CFB87C",
    "green": "#2D6A4F",
    "ink": "#1A1A2E",
    "body": "#3A4658",
    "muted": "#6B778A",
    "light_bg": "#F0F4FA",
    "divider": "#D0D8E8",
    "white": "#FFFFFF",
    "black": "#000000",
}


def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ── Theme presets ──────────────────────────────────────────────────
THEME_PRESETS = {
    "petro-blue": {
        "primary": "blue",
        "accent": "dark_blue",
        "bg": "white",
        "text": "ink",
        "title_bar": "blue",
    },
    "nepu-red": {
        "primary": "red",
        "accent": "dark_blue",
        "bg": "white",
        "text": "ink",
        "title_bar": "red",
    },
    "petro-gold": {
        "primary": "dark_blue",
        "accent": "gold",
        "bg": "dark_blue",
        "text": "white",
        "body": "divider",
        "muted": "gold",
        "light_bg": "ink",
        "title_bar": "dark_blue",
    },
    "oilfield-green": {
        "primary": "green",
        "accent": "gold",
        "bg": "white",
        "text": "ink",
        "title_bar": "green",
    },
    "clean-white": {
        "primary": "blue",
        "accent": "muted",
        "bg": "white",
        "text": "ink",
        "title_bar": "ink",
    },
}


class NepuSlideBuilder:
    """Build a NEPU-styled PPTX deck with high-level slide patterns.

    Defaults to 16:9 widescreen (13.333" × 7.5").
    """

    def __init__(self, theme: str = "petro-blue", slide_width: float = 13.333,
                 slide_height: float = 7.5):
        if theme not in THEME_PRESETS:
            raise ValueError(f"Unknown theme '{theme}'. Choose from: {list(THEME_PRESETS)}")
        self.theme_name = theme
        self.colors = THEME_PRESETS[theme]
        self.prs = Presentation()
        self.prs.slide_width = Inches(slide_width)
        self.prs.slide_height = Inches(slide_height)
        self.sw = slide_width   # inches
        self.sh = slide_height  # inches
        self._slide_count = 0
        # Use blank layout by default
        self._blank_layout = self.prs.slide_layouts[6]  # blank

    # ── internal helpers ───────────────────────────────────────────

    def _add_slide(self) -> Any:
        slide = self.prs.slides.add_slide(self._blank_layout)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = _rgb(self._c("bg"))
        self._slide_count += 1
        return slide

    def _c(self, key: str) -> str:
        """Resolve color key → hex string."""
        color_key = self.colors.get(key, key)
        return PALETTE.get(color_key, PALETTE["ink"])

    @staticmethod
    def _require_file(path: str | Path, label: str) -> str:
        resolved = Path(path).expanduser().resolve()
        if not resolved.is_file():
            raise FileNotFoundError(f"{label} does not exist: {resolved}")
        return str(resolved)

    @staticmethod
    def _add_textbox(slide, left: float, top: float, width: float, height: float,
                     text: str = "", font_size: float = 16, bold: bool = False,
                     color: str = "ink", alignment: int = PP_ALIGN.LEFT,
                     font_name: str = "Microsoft YaHei") -> Any:
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = _rgb(color)
        p.font.name = font_name
        p.alignment = alignment
        return txBox

    @staticmethod
    def _add_rect(slide, left: float, top: float, width: float, height: float,
                  fill_color: str = "blue") -> Any:
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(left), Inches(top), Inches(width), Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill_color)
        shape.line.fill.background()
        return shape

    @staticmethod
    def _add_picture_contain(
        slide,
        path: str,
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> Any:
        """Add a picture without distortion, centered inside an inch-based box."""
        box_left = Inches(left)
        box_top = Inches(top)
        box_width = Inches(width)
        box_height = Inches(height)
        picture = slide.shapes.add_picture(path, box_left, box_top)
        scale = min(box_width / picture.width, box_height / picture.height)
        picture.width = int(picture.width * scale)
        picture.height = int(picture.height * scale)
        picture.left = int(box_left + (box_width - picture.width) / 2)
        picture.top = int(box_top + (box_height - picture.height) / 2)
        return picture

    # ── slide patterns ─────────────────────────────────────────────

    def add_cover(self, title: str, *, subtitle: str = "",
                  logo_path: str | None = None) -> None:
        """Full-bleed cover slide with primary-color background."""
        logo = self._require_file(logo_path, "Logo") if logo_path else None
        slide = self._add_slide()

        # Full-slide background rect
        cover_bg = self._c("bg") if self.theme_name == "clean-white" else self._c("primary")
        self._add_rect(slide, 0, 0, self.sw, self.sh, fill_color=cover_bg)

        # Logo (top-right, small)
        if logo:
            self._add_picture_contain(slide, logo, self.sw - 2.2, 0.4, 1.8, 0.7)

        # Title — centered, bold
        title_y = self.sh * 0.35
        self._add_textbox(
            slide, 1.0, title_y, self.sw - 2.0, 1.5,
            text=title, font_size=36, bold=True,
            color=self._c("text") if self.theme_name == "clean-white" else self._c("white"),
            alignment=PP_ALIGN.CENTER,
        )

        # Subtitle
        if subtitle:
            subtitle_color = (
                self._c("text")
                if self.theme_name == "clean-white"
                else self._c("accent")
                if self.theme_name == "petro-gold"
                else self._c("white")
            )
            self._add_textbox(
                slide, 1.5, title_y + 1.5, self.sw - 3.0, 0.8,
                text=subtitle, font_size=18, bold=False,
                color=subtitle_color,
                alignment=PP_ALIGN.CENTER,
            )

        # Bottom decorative line
        self._add_rect(
            slide, self.sw * 0.3, self.sh - 0.15, self.sw * 0.4, 0.04,
            fill_color=self._c("accent"),
        )

    def add_section(self, title: str, *, subtitle: str = "") -> None:
        """Section divider slide — dark background with centered title."""
        slide = self._add_slide()

        # Full-slide dark background
        bg_color = (
            self._c("primary")
            if self.theme_name in {"petro-gold", "oilfield-green"}
            else self._c("accent")
        )
        self._add_rect(slide, 0, 0, self.sw, self.sh, fill_color=bg_color)

        # Left accent bar
        self._add_rect(slide, 0.6, self.sh * 0.3, 0.06, self.sh * 0.4,
                       fill_color=self._c("gold"))

        # Section number and title
        self._add_textbox(
            slide, 1.2, self.sh * 0.33, self.sw - 2.4, 1.0,
            text=f"0{self._slide_count}" if self._slide_count < 10 else str(self._slide_count),
            font_size=14, bold=True,
            color=self._c("muted"),
        )
        self._add_textbox(
            slide, 1.2, self.sh * 0.4, self.sw - 2.4, 1.2,
            text=title, font_size=30, bold=True,
            color=self._c("white"),
        )
        if subtitle:
            self._add_textbox(
                slide, 1.2, self.sh * 0.55, self.sw - 2.4, 0.6,
                text=subtitle, font_size=14, bold=False,
                color=self._c("muted"),
            )

    def add_claim_slide(self, claim: str, *,
                        bullets: list[str] | None = None) -> None:
        """Claim-led slide: one bold claim sentence + optional support bullets."""
        slide = self._add_slide()

        # Top color bar
        self._add_rect(slide, 0, 0, self.sw, 0.08, fill_color=self._c("primary"))

        # Main claim
        self._add_textbox(
            slide, 1.0, 1.5, self.sw - 2.0, 2.0,
            text=claim, font_size=28, bold=True,
            color=self._c("text"),
            alignment=PP_ALIGN.LEFT,
        )

        # Bullets
        if bullets:
            bullet_text = "\n".join(f"• {b}" for b in bullets)
            self._add_textbox(
                slide, 1.0, 4.0, self.sw - 2.0, 2.5,
                text=bullet_text, font_size=14, bold=False,
                color=self._c("body"),
            )

        # Bottom line
        self._add_rect(slide, 1.0, self.sh - 0.6, self.sw - 2.0, 0.02,
                       fill_color=self._c("divider"))

    def add_figure_slide(self, figure_path: str, *, claim: str = "",
                         layout: str = "asymmetric-70-30",
                         caption: str = "",
                         notes: str = "") -> None:
        """Figure-dominant slide with optional claim/caption sidebar.

        Layouts:
          - "full": figure fills entire slide
          - "asymmetric-70-30": figure left 70%, text right 30%
          - "top-title": title bar at top, figure below
        """
        if layout not in {"full", "asymmetric-70-30", "top-title"}:
            raise ValueError(f"Unknown figure layout: {layout}")
        figure = self._require_file(figure_path, "Figure")
        slide = self._add_slide()

        if layout == "full":
            self._add_picture_contain(slide, figure, 0, 0, self.sw, self.sh)

        elif layout == "asymmetric-70-30":
            # Left: figure
            fig_w = self.sw * 0.68
            self._add_picture_contain(
                slide, figure, 0.3, 0.5, fig_w - 0.3, self.sh - 1.2
            )
            # Right: claim column
            if claim:
                self._add_textbox(
                    slide, fig_w + 0.3, 1.5, self.sw - fig_w - 0.8, 2.0,
                    text=claim, font_size=18, bold=True,
                    color=self._c("text"),
                )
            if caption:
                self._add_textbox(
                    slide, fig_w + 0.3, self.sh - 2.0, self.sw - fig_w - 0.8, 1.2,
                    text=caption, font_size=9, bold=False,
                    color=self._c("muted"),
                )

        elif layout == "top-title":
            if claim:
                self._add_textbox(
                    slide, 0.8, 0.3, self.sw - 1.6, 0.7,
                    text=claim, font_size=18, bold=True,
                    color=self._c("text"),
                )
                # Underline
                self._add_rect(slide, 0.8, 0.95, 1.5, 0.03,
                               fill_color=self._c("primary"))
            self._add_picture_contain(
                slide, figure, 0.5, 1.2, self.sw - 1.0, self.sh - 1.8
            )

        # Speaker notes
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def add_matrix_slide(self, headers: list[str], rows: list[list[str]], *,
                         title: str = "",
                         header_color: str | None = None) -> None:
        """Comparison/matrix slide with a simple table."""
        if not headers:
            raise ValueError("Matrix slide requires at least one header")
        for index, row in enumerate(rows, start=1):
            if len(row) != len(headers):
                raise ValueError(
                    f"Matrix row {index} has {len(row)} cells; expected {len(headers)}"
                )
        slide = self._add_slide()

        if header_color is None:
            header_color = self._c("primary")

        # Title
        if title:
            self._add_textbox(
                slide, 0.8, 0.3, self.sw - 1.6, 0.7,
                text=title, font_size=22, bold=True,
                color=self._c("text"),
            )

        # Table dimensions
        n_rows = len(rows) + 1  # +1 for header
        n_cols = len(headers)
        table_left = 0.8
        table_top = 1.2 if title else 0.5
        table_w = self.sw - 1.6
        table_h = min(self.sh - table_top - 0.5, max(1.4, n_rows * 0.65))

        table = slide.shapes.add_table(n_rows, n_cols,
                                       Inches(table_left), Inches(table_top),
                                       Inches(table_w), Inches(table_h)).table

        # Column widths
        col_w = table_w / n_cols
        for i in range(n_cols):
            table.columns[i].width = Inches(col_w)

        # Header row
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = _rgb(self._c("white"))
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(header_color)

        # Data rows
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.text = str(val)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(12)
                    p.font.color.rgb = _rgb(self._c("text"))
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(
                    self._c("light_bg") if i % 2 == 1 else self._c("bg")
                )

    def add_bullet_slide(self, title: str, bullets: list[str], *,
                         two_column: bool = False) -> None:
        """Standard text slide with title and bullet points."""
        slide = self._add_slide()

        # Top color bar
        self._add_rect(slide, 0, 0, self.sw, 0.06, fill_color=self._c("primary"))

        # Title
        self._add_textbox(
            slide, 0.8, 0.3, self.sw - 1.6, 0.7,
            text=title, font_size=24, bold=True, color=self._c("text"),
        )
        # Underline
        self._add_rect(slide, 0.8, 0.95, 1.2, 0.03, fill_color=self._c("primary"))

        if two_column:
            mid = len(bullets) // 2 + len(bullets) % 2
            left = bullets[:mid]
            right = bullets[mid:]
            for side, offset, lst in [("L", 0, left), ("R", 0.48, right)]:
                text = "\n".join(f"• {b}" for b in lst)
                self._add_textbox(
                    slide, 0.8 + self.sw * offset, 1.4, self.sw * 0.43 - 0.8, self.sh - 2.0,
                    text=text, font_size=15, bold=False, color=self._c("body"),
                )
        else:
            bullet_text = "\n".join(f"• {b}" for b in bullets)
            self._add_textbox(
                slide, 0.8, 1.4, self.sw - 1.6, self.sh - 2.0,
                text=bullet_text, font_size=16, bold=False, color=self._c("body"),
            )

    def add_conclusion(self, takeaways: list[str], *,
                       thank_you: str = "谢谢！") -> None:
        """Closing slide with key takeaways and thank-you line."""
        slide = self._add_slide()

        # Full background
        self._add_rect(slide, 0, 0, self.sw, self.sh, fill_color=self._c("primary"))

        # Takeaways
        if takeaways:
            text = "\n".join(f"{i+1}. {t}" for i, t in enumerate(takeaways))
            self._add_textbox(
                slide, 1.5, self.sh * 0.2, self.sw - 3.0, self.sh * 0.5,
                text=text, font_size=18, bold=False,
                color=self._c("white"),
                alignment=PP_ALIGN.LEFT,
            )

        # Thank you
        self._add_textbox(
            slide, 1.5, self.sh * 0.70, self.sw - 3.0, 1.0,
            text=thank_you, font_size=30, bold=True,
            color=self._c("white"),
            alignment=PP_ALIGN.CENTER,
        )

        # Bottom accent
        self._add_rect(slide, self.sw * 0.35, self.sh - 0.12, self.sw * 0.3, 0.04,
                       fill_color=self._c("accent"))

    def add_image_grid_slide(self, image_paths: list[str], *,
                             title: str = "", cols: int = 3,
                             captions: list[str] | None = None) -> None:
        """Grid of images with optional title and captions."""
        if cols < 1:
            raise ValueError("Image grid columns must be at least 1")
        if not image_paths:
            raise ValueError("Image grid requires at least one image")
        images = [self._require_file(path, f"Image {index}") for index, path in enumerate(image_paths, start=1)]
        slide = self._add_slide()

        if title:
            self._add_textbox(
                slide, 0.8, 0.3, self.sw - 1.6, 0.6,
                text=title, font_size=18, bold=True, color=self._c("text"),
            )

        rows = (len(image_paths) + cols - 1) // cols
        cell_w = (self.sw - 1.0) / cols
        cell_h = (self.sh - 1.8) / max(rows, 1)
        top_offset = 1.2 if title else 0.4

        for i, img_path in enumerate(images):
            r, c = divmod(i, cols)
            left = 0.5 + c * cell_w + 0.1
            top = top_offset + r * cell_h + 0.1
            self._add_picture_contain(
                slide, img_path, left, top, cell_w - 0.2, cell_h - 0.35
            )
            if captions and i < len(captions):
                self._add_textbox(
                    slide, left, top + cell_h - 0.35, cell_w - 0.2, 0.3,
                    text=captions[i], font_size=9, bold=False,
                    color=self._c("muted"), alignment=PP_ALIGN.CENTER,
                )

    def add_blank_slide(self) -> Any:
        """Return a blank slide for fully custom layout."""
        return self._add_slide()

    def set_speaker_notes(self, slide_number: int, notes: str) -> None:
        """Set speaker notes on a 1-based slide number."""
        if slide_number < 1 or slide_number > self._slide_count:
            raise IndexError(
                f"Slide number {slide_number} is outside 1-{self._slide_count}"
            )
        self.prs.slides[slide_number - 1].notes_slide.notes_text_frame.text = notes

    # ── output ──────────────────────────────────────────────────────

    def save(self, path: str | Path, *, overwrite: bool = False) -> str:
        """Save the presentation and return the absolute path.

        Refuses to overwrite by default to protect user-edited and delivered files.
        """
        output = Path(path).resolve()
        if output.exists() and not overwrite:
            raise FileExistsError(
                f"Refusing to overwrite existing presentation: {output}"
            )
        output.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output))
        return str(output)

    @property
    def slide_count(self) -> int:
        return self._slide_count
